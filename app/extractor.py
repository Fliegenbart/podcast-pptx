"""
Extrahiert Inhalte aus PPTX-Dateien, inklusive Text und Bilder.
"""
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
import logging
import io

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

logger = logging.getLogger(__name__)


@dataclass
class SlideImage:
    """Ein Bild aus einer Folie."""
    image_bytes: bytes
    content_type: str
    width: Optional[int] = None
    height: Optional[int] = None
    alt_text: Optional[str] = None
    position_hint: str = "center"  # left, center, right, full
    description: Optional[str] = None  # Wird später durch LLaVA gefüllt


@dataclass
class SlideContent:
    """Inhalte einer einzelnen Folie."""
    slide_number: int
    title: Optional[str] = None
    body_text: list[str] = field(default_factory=list)
    notes: Optional[str] = None
    images: list[SlideImage] = field(default_factory=list)
    section_name: Optional[str] = None
    is_section_start: bool = False

    def has_content(self) -> bool:
        """Prüft ob die Folie relevanten Inhalt hat."""
        return bool(self.title or self.body_text or self.images)

    def to_text(self) -> str:
        """Konvertiert Folie zu lesbarem Text."""
        parts = []
        if self.title:
            parts.append(f"**{self.title}**")
        if self.body_text:
            parts.extend(self.body_text)
        if self.notes:
            parts.append(f"[Notizen: {self.notes}]")

        # Bildbeschreibungen einfügen
        for img in self.images:
            if img.description:
                parts.append(f"[Bild: {img.description}]")
            elif img.alt_text:
                parts.append(f"[Bild: {img.alt_text}]")

        return "\n".join(parts)


@dataclass
class PresentationContent:
    """Gesamter Präsentationsinhalt."""
    title: Optional[str] = None
    slides: list[SlideContent] = field(default_factory=list)
    total_slides: int = 0
    sections: list[str] = field(default_factory=list)

    def to_prompt_text(self) -> str:
        """Formatiert Inhalte für LLM-Prompt."""
        parts = []

        if self.title:
            parts.append(f"# {self.title}\n")

        current_section = None
        for slide in self.slides:
            if slide.section_name and slide.section_name != current_section:
                current_section = slide.section_name
                parts.append(f"\n## Abschnitt: {current_section}\n")

            parts.append(f"\n### Folie {slide.slide_number}")
            parts.append(slide.to_text())

        return "\n".join(parts)

    def get_slide_titles(self) -> list[tuple[int, str]]:
        """Gibt Liste von (Foliennummer, Titel) zurück."""
        return [
            (slide.slide_number, slide.title or f"Folie {slide.slide_number}")
            for slide in self.slides
        ]


def _extract_text_from_shape(shape) -> list[str]:
    """Extrahiert Text aus einem Shape."""
    texts = []

    if hasattr(shape, "text") and shape.text.strip():
        texts.append(shape.text.strip())

    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                texts.append(text)

    return texts


def _extract_image_from_shape(shape, slide_width: int) -> Optional[SlideImage]:
    """Extrahiert ein Bild aus einem Shape."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_bytes = image.blob
            content_type = image.content_type

            # Position bestimmen
            left = shape.left
            if left < slide_width * 0.33:
                position = "left"
            elif left > slide_width * 0.66:
                position = "right"
            else:
                position = "center"

            # Größe ermitteln
            try:
                img = Image.open(io.BytesIO(image_bytes))
                width, height = img.size
            except Exception:
                width, height = None, None

            # Alt-Text wenn vorhanden
            alt_text = None
            if hasattr(shape, "_element"):
                desc_elem = shape._element.find(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr"
                )
                if desc_elem is not None:
                    alt_text = desc_elem.get("descr")

            return SlideImage(
                image_bytes=image_bytes,
                content_type=content_type,
                width=width,
                height=height,
                alt_text=alt_text,
                position_hint=position
            )
    except Exception as e:
        logger.warning(f"Fehler beim Extrahieren eines Bildes: {e}")

    return None


def _get_slide_notes(slide) -> Optional[str]:
    """Extrahiert Notizen einer Folie."""
    try:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                return notes_text
    except Exception:
        pass
    return None


def extract_pptx(pptx_path: Path) -> PresentationContent:
    """
    Extrahiert alle Inhalte aus einer PPTX-Datei.

    Args:
        pptx_path: Pfad zur PPTX-Datei

    Returns:
        PresentationContent mit allen extrahierten Inhalten
    """
    logger.info(f"Extrahiere PPTX: {pptx_path}")

    prs = Presentation(str(pptx_path))

    # Präsentationstitel (erste Folie)
    presentation_title = None

    # Slide-Breite für Position-Berechnung
    slide_width = prs.slide_width

    # Sections sammeln (falls vorhanden)
    sections = []
    section_map = {}  # slide_index -> section_name

    # PowerPoint-Sections extrahieren (falls verfügbar)
    try:
        for section in prs.sections:
            sections.append(section.name)
            for slide_idx in section.slide_idxs:
                section_map[slide_idx] = section.name
    except AttributeError:
        # Keine Sections vorhanden
        pass

    slides = []
    previous_section = None

    for idx, slide in enumerate(prs.slides):
        slide_number = idx + 1

        # Section-Info
        section_name = section_map.get(idx)
        is_section_start = section_name is not None and section_name != previous_section
        previous_section = section_name

        slide_content = SlideContent(
            slide_number=slide_number,
            section_name=section_name,
            is_section_start=is_section_start
        )

        # Durch alle Shapes iterieren
        for shape in slide.shapes:
            # Titel
            if shape.is_placeholder and hasattr(shape, "placeholder_format"):
                if shape.placeholder_format.type == 1:  # Title
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.title = shape.text.strip()
                        if idx == 0 and not presentation_title:
                            presentation_title = shape.text.strip()
                        continue

            # Bilder extrahieren
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = _extract_image_from_shape(shape, slide_width)
                if img:
                    slide_content.images.append(img)
                continue

            # Gruppenformen durchsuchen
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img = _extract_image_from_shape(sub_shape, slide_width)
                        if img:
                            slide_content.images.append(img)
                    else:
                        texts = _extract_text_from_shape(sub_shape)
                        slide_content.body_text.extend(texts)
                continue

            # Sonstiger Text
            texts = _extract_text_from_shape(shape)
            # Titel nicht doppelt hinzufügen
            for text in texts:
                if text != slide_content.title and text not in slide_content.body_text:
                    slide_content.body_text.append(text)

        # Notizen
        slide_content.notes = _get_slide_notes(slide)

        # Nur Folien mit Inhalt hinzufügen
        if slide_content.has_content():
            slides.append(slide_content)

    content = PresentationContent(
        title=presentation_title,
        slides=slides,
        total_slides=len(prs.slides),
        sections=sections
    )

    logger.info(
        f"Extrahiert: {len(slides)} Folien mit Inhalt, "
        f"{sum(len(s.images) for s in slides)} Bilder"
    )

    return content
